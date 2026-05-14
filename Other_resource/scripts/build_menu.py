#!/usr/bin/env python3
"""Build swiss-bakery-menu.pptx — A4 portrait, editable template for Canva import.

Run: /tmp/swiss-venv/bin/python scripts/build_menu.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
LOGO = ROOT / "swiss_logo.png"
OUT = ROOT / "swiss-bakery-menu.pptx"

NAVY = RGBColor(0x1A, 0x2A, 0x4F)
GOLD = RGBColor(0xC8, 0x9A, 0x2B)
CREAM = RGBColor(0xF6, 0xEF, 0xE0)
RED = RGBColor(0xC8, 0x10, 0x2E)
DIM = RGBColor(0x8A, 0x7A, 0x5A)

HEADER_FONT = "Playfair Display"
BODY_FONT = "Montserrat"

SAVOURY = [
    ("Chicken · Paneer · Beef Patties", "$3.50"),
    ("Chicken / Tuna Sandwich", "$6.00"),
    ("Small Burger", "$4.00"),
    ("Prawn Roll", "$4.00"),
    ("Chicken Roll", "$4.00"),
    ("Beef Roll", "$4.00"),
    ("Samosa", "$3.00"),
    ("Singara (Veg / Liver)", "$2.80"),
    ("Piyaju · Dal Puri", "$2.00"),
    ("Mughlai Paratha", "$10.00"),
    ("Popcorn Chicken", "ask staff"),
]

BAKERY = [
    ("Swiss Short Bread (2 pcs)", "$2.20"),
    ("Salty Cookies · Nimki", "ask staff"),
    ("Butter Bun", "ask staff"),
    ("Cream Roll", "ask staff"),
    ("Kaza", "ask staff"),
    ("Cakes — Fruit & Others", "ask staff"),
    ("Patishapta", "ask staff"),
    ("Pakon", "ask staff"),
    ("Narikel Puli (Coconut)", "ask staff"),
]

SWEETS = [
    ("Laddu — 3 varieties", "$23 / kg"),
    ("Monsor", "$30 / kg"),
    ("Kodom", "$32 / kg"),
    ("Balu Shai", "$25 / kg"),
    ("Aflaton", "$30 / kg"),
    ("Pera", "$35 / kg"),
    ("Bondia", "ask staff"),
    ("Halwa — Chana Dal & Hazelnut", "$20 / kg"),
    ("Chana Mukhi", "$35 / kg"),
    ("Rasgulla", "$25 / kg"),
    ("Cham Cham", "$26 / kg"),
    ("Rosmalai", "$30 / kg"),
    ("Hafsa", "$35 / kg"),
]

BEVERAGES = [
    ("Malai Cha — small cup", "$2.80"),
    ("Poro Roti", "ask staff"),
    ("Buttered Corn", "ask staff"),
]


def style_run(run, font, size, color, bold=False, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=10, color=NAVY,
             bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], font, size, color, bold, italic)
    return tb


def add_line(slide, x1, y1, x2, y2, color=GOLD, weight=0.6):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Mm(x1), Mm(y1), Mm(x2), Mm(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    return conn


def add_rect(slide, x, y, w, h, fill=None, stroke=None, stroke_w=0.6):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(y), Mm(w), Mm(h))
    if fill is None:
        r.fill.background()
    else:
        r.fill.solid()
        r.fill.fore_color.rgb = fill
    if stroke is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = stroke
        r.line.width = Pt(stroke_w)
    return r


def section_header(slide, y, title, subtitle):
    add_line(slide, 22, y + 3.6, 48, y + 3.6, color=GOLD, weight=0.5)
    add_text(slide, 50, y, 110, 7, title,
             font=HEADER_FONT, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 162, y + 3.6, 188, y + 3.6, color=GOLD, weight=0.5)
    if subtitle:
        add_text(slide, 15, y + 7.5, 180, 4.5, subtitle.upper(),
                 font=BODY_FONT, size=7.5, color=GOLD, italic=True,
                 align=PP_ALIGN.CENTER)


def item_rows(slide, items, start_y, row_h=5.5, two_cols=True):
    col_w = 87
    name_w = col_w - 24
    price_w = 24
    gap = 4
    left_x = 13
    right_x = left_x + col_w + gap
    if two_cols:
        mid = (len(items) + 1) // 2
        cols = [(items[:mid], left_x), (items[mid:], right_x)]
    else:
        cols = [(items, left_x)]
        col_w = 184
        name_w = col_w - 24

    for col_items, x in cols:
        for i, (name, price) in enumerate(col_items):
            yy = start_y + i * row_h
            add_text(slide, x, yy, name_w, row_h, name, size=9, color=NAVY)
            is_price = price.startswith("$")
            add_text(slide, x + col_w - price_w, yy, price_w, row_h, price,
                     size=9,
                     color=GOLD if is_price else DIM,
                     bold=is_price,
                     italic=not is_price,
                     align=PP_ALIGN.RIGHT)
    rows = max(len(col_items) for col_items, _ in cols)
    return start_y + rows * row_h


def build():
    prs = Presentation()
    prs.slide_width = Mm(210)
    prs.slide_height = Mm(297)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 210, 297, fill=CREAM)
    add_rect(slide, 7, 7, 196, 283, stroke=GOLD, stroke_w=1.1)
    add_rect(slide, 9, 9, 192, 279, stroke=GOLD, stroke_w=0.3)

    if LOGO.exists():
        logo_w = 42
        slide.shapes.add_picture(str(LOGO),
                                 Mm((210 - logo_w) / 2), Mm(14),
                                 width=Mm(logo_w))

    add_text(slide, 15, 59, 180, 7,
             "Handcrafted daily · Swiss soul, Bengali heart",
             font=HEADER_FONT, size=12, color=NAVY, italic=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 15, 67, 180, 4.5,
             "ALL PRICES AUD  ·  DINE-IN & TAKEAWAY",
             font=BODY_FONT, size=7.5, color=GOLD, bold=True,
             align=PP_ALIGN.CENTER)

    add_line(slide, 25, 75, 185, 75, color=GOLD, weight=0.9)

    y = 81
    section_header(slide, y, "SAVOURY BITES", "per piece")
    y = item_rows(slide, SAVOURY, y + 13) + 4

    section_header(slide, y, "BAKERY & BREADS", "freshly baked in-house")
    y = item_rows(slide, BAKERY, y + 13) + 4

    section_header(slide, y, "TRADITIONAL SWEETS", "priced by weight")
    y = item_rows(slide, SWEETS, y + 13) + 4

    section_header(slide, y, "BEVERAGES & SIDES", "")
    y = item_rows(slide, BEVERAGES, y + 10, two_cols=False)

    add_line(slide, 25, 278, 185, 278, color=GOLD, weight=0.4)
    add_text(slide, 15, 280, 180, 5.5,
             "Baked fresh every morning",
             font=HEADER_FONT, size=10, color=NAVY, italic=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 15, 286, 180, 4,
             "SWISS  ✚  BAKERY   ·   AUSTRALIA",
             font=BODY_FONT, size=7.5, color=GOLD, bold=True,
             align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
